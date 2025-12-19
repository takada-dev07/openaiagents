from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from app.core.config import Settings, get_settings, new_trace_id
from app.core.logging import get_logger
from app.models.schemas import ExplainedSlide, PptxSlideSpec
from app.utils.convert import ConvertError, pptx_to_pngs


@dataclass(frozen=True)
class PptxRenderResult:
    trace_id: str
    pptx_path: str
    artifacts: list[str]


@dataclass(frozen=True)
class PptxExplainResult:
    trace_id: str
    slides: list[ExplainedSlide]
    images: list[str]
    warnings: list[str]


class PptxService:
    def __init__(self, *, settings: Settings | None = None, template_path: Path | None = None) -> None:
        self.settings = settings or get_settings()
        self.log = get_logger(__name__)
        self.template_path = template_path or (Path("templates") / "template.pptx")

    async def render_deck(self, *, title: str, slides: list[PptxSlideSpec]) -> PptxRenderResult:
        trace_id = new_trace_id()
        out_path = self.settings.artifact_dir / f"deck-{trace_id}.pptx"

        prs = self._load_template_or_blank()

        # Title slide
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            if slide.shapes.title:
                slide.shapes.title.text = self._truncate(title, 60)
        except Exception:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = self._truncate(title, 60)

        for spec in slides:
            self._add_content_slide(prs, spec)

        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(out_path)

        meta_path = self.settings.artifact_dir / "traces" / f"{trace_id}-pptx.json"
        meta_path.write_text(
            json.dumps(
                {
                    "trace_id": trace_id,
                    "pptx_path": str(out_path),
                    "title": title,
                    "slide_count": len(prs.slides),
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )

        artifacts = [str(out_path), str(meta_path)]
        self.log.info("Rendered PPTX", extra={"trace_id": trace_id, "pptx_path": str(out_path)})
        return PptxRenderResult(trace_id=trace_id, pptx_path=str(out_path), artifacts=artifacts)

    async def explain_deck(self, *, pptx_path: str) -> PptxExplainResult:
        trace_id = new_trace_id()
        warnings: list[str] = []

        path = Path(pptx_path)
        prs = Presentation(path)
        explained: list[ExplainedSlide] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                t = (shape.text or "").strip()
                if t:
                    texts.append(t)
            all_text = "\n".join(texts).strip()
            notes = ""
            try:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            except Exception:
                notes = ""

            explain = self._rule_based_explain(all_text)
            explained.append(ExplainedSlide(slide=i, text=all_text, notes=notes, explain=explain))

        images: list[str] = []
        out_dir = self.settings.artifact_dir / "slides" / trace_id
        try:
            conv = pptx_to_pngs(pptx_path=path, out_dir=out_dir)
            images = [str(p) for p in conv.png_paths]
        except FileNotFoundError as e:
            warnings.append(f"変換ツールが見つかりません: {e}")
        except ConvertError as e:
            warnings.append(f"PPTX→画像変換に失敗しました: {e}")
        except Exception as e:
            warnings.append(f"PPTX→画像変換で予期しないエラー: {e}")

        return PptxExplainResult(trace_id=trace_id, slides=explained, images=images, warnings=warnings)

    def _load_template_or_blank(self) -> Presentation:
        try:
            if self.template_path.exists():
                return Presentation(str(self.template_path))
        except Exception as e:
            self.log.warning("Failed to load template, falling back to blank", extra={"error": str(e)})
        return Presentation()

    def _add_content_slide(self, prs: Presentation, spec: PptxSlideSpec) -> None:
        # Default “Title and Content”
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = self._truncate(spec.heading, 60)

        body = None
        # best-effort: placeholder[1] often is body
        try:
            body = slide.shapes.placeholders[1].text_frame
        except Exception:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape != slide.shapes.title:
                    body = shape.text_frame
                    break

        if body is not None:
            body.clear()
            for line in self._truncate_bullets(spec.bullets):
                p = body.add_paragraph()
                p.text = line
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(18)

        if spec.image_path:
            self._add_image_contain(slide, Path(spec.image_path))

    def _add_image_contain(self, slide, image_path: Path) -> None:
        # Right side box
        left = Inches(6.5)
        top = Inches(1.6)
        width = Inches(3.0)
        height = Inches(3.2)

        try:
            with Image.open(image_path) as img:
                px_w, px_h = img.size
        except Exception:
            slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
            return

        # Calculate contain scaling in EMU units
        box_w = width
        box_h = height
        img_ratio = px_w / px_h if px_h else 1.0
        box_ratio = box_w / box_h

        if img_ratio >= box_ratio:
            # fit to width
            w = box_w
            h = box_w / img_ratio
        else:
            h = box_h
            w = box_h * img_ratio

        x = left + (box_w - w) / 2
        y = top + (box_h - h) / 2
        slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)

    def _truncate_bullets(self, bullets: list[str], *, max_lines: int = 6, max_chars: int = 60) -> list[str]:
        out: list[str] = []
        for b in bullets[:max_lines]:
            out.append(self._truncate(b, max_chars))
        if len(bullets) > max_lines:
            out.append("…")
        return out

    def _truncate(self, s: str, max_chars: int) -> str:
        s = (s or "").strip()
        if len(s) <= max_chars:
            return s
        return s[: max(0, max_chars - 1)] + "…"

    def _rule_based_explain(self, text: str) -> str:
        if not text:
            return "このスライドはテキスト要素が少ないため、図や画像の説明が中心になりそうです。"
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        heading = lines[0] if lines else ""
        bullets = lines[1:4]
        if bullets:
            return f"見出し「{heading}」の要点を、箇条書き（{len(bullets)}件）で説明しています。"
        return f"見出し「{heading}」について短く説明するスライドです。"


